#!/usr/bin/env python3
"""Build formal UCS503P W4 pitch deck (DeadlineDesk)."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml

# Academic palette (TIET-adjacent: deep navy + amber accent)
NAVY = RGBColor(0x0B, 0x1F, 0x3A)
NAVY2 = RGBColor(0x14, 0x32, 0x56)
AMBER = RGBColor(0xC4, 0x7E, 0x1A)
AMBER_LT = RGBColor(0xF5, 0xEB, 0xD7)
INK = RGBColor(0x1C, 0x24, 0x33)
MUTED = RGBColor(0x4E, 0x5A, 0x6A)
LINE = RGBColor(0xC9, 0xD0, 0xD8)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
PAPER = RGBColor(0xF7, 0xF8, 0xFA)
GREEN = RGBColor(0x1B, 0x6B, 0x4A)


def font(run, size=16, bold=False, color=INK, name="Calibri"):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = name


def rect(slide, l, t, w, h, fill, line=None):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
    return sh


def txt(slide, l, t, w, h, lines, size=16, bold=False, color=INK, align=PP_ALIGN.LEFT, name="Calibri"):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(4)
        r = p.add_run()
        r.text = line
        font(r, size, bold, color, name)
    return box


def footer(slide, prs, n, total, label="UCS503P · Project Proposal · W4"):
    rect(slide, 0, prs.slide_height - Inches(0.42), prs.slide_width, Inches(0.42), NAVY)
    txt(
        slide,
        Inches(0.45),
        prs.slide_height - Inches(0.38),
        Inches(10),
        Inches(0.3),
        [label],
        11,
        False,
        RGBColor(0xC5, 0xD0, 0xDC),
    )
    txt(
        slide,
        prs.slide_width - Inches(1.2),
        prs.slide_height - Inches(0.38),
        Inches(0.9),
        Inches(0.3),
        [f"{n} / {total}"],
        11,
        True,
        AMBER,
        PP_ALIGN.RIGHT,
    )


def header(slide, prs, section, title):
    rect(slide, 0, 0, prs.slide_width, Inches(1.15), NAVY)
    rect(slide, 0, Inches(1.15), prs.slide_width, Inches(0.06), AMBER)
    txt(slide, Inches(0.5), Inches(0.22), Inches(12), Inches(0.28), [section.upper()], 11, True, AMBER)
    txt(slide, Inches(0.5), Inches(0.48), Inches(12), Inches(0.55), [title], 26, True, WHITE, name="Georgia")


def table_slide(slide, prs, rows, col_w, left=Inches(0.5), top=Inches(1.55)):
    cols = len(rows[0])
    row_h = Inches(0.42)
    y = top
    for ri, row in enumerate(rows):
        x = left
        bg = NAVY if ri == 0 else (AMBER_LT if ri % 2 == 0 else WHITE)
        fg = WHITE if ri == 0 else INK
        for ci, cell in enumerate(row):
            w = col_w[ci]
            rect(slide, x, y, w, row_h, bg, LINE)
            txt(slide, x + Inches(0.1), y + Inches(0.08), w - Inches(0.15), row_h, [cell], 12, ri == 0, fg)
            x += w
        y += row_h


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    N = 10

    # ----- 1 Title -----
    s = prs.slides.add_slide(blank)
    rect(s, 0, 0, prs.slide_width, prs.slide_height, NAVY)
    rect(s, 0, Inches(5.9), prs.slide_width, Inches(1.6), NAVY2)
    rect(s, Inches(0.5), Inches(1.7), Inches(1.2), Inches(0.08), AMBER)
    txt(s, Inches(0.5), Inches(0.55), Inches(12), Inches(0.35),
        ["THAPAR INSTITUTE OF ENGINEERING AND TECHNOLOGY  ·  UCS503P"], 12, True, AMBER)
    txt(s, Inches(0.5), Inches(1.95), Inches(12), Inches(1.4),
        ["DeadlineDesk"], 44, True, WHITE, name="Georgia")
    txt(s, Inches(0.5), Inches(3.35), Inches(11.5), Inches(1.2),
        [
            "A Campus Web System for Placement Rounds",
            "and Academic Assignment Deadlines",
        ],
        22, False, RGBColor(0xD7, 0xDE, 0xE7), name="Georgia")
    txt(s, Inches(0.5), Inches(6.15), Inches(12), Inches(1.0),
        [
            "Project Proposal Presentation  ·  Week 4",
            "Aryan Gupta (1024030764)  ·  Aksh Goyal (1024030766)  ·  Naveen Bansal (1024030767)",
            "Submitted to: Dr. Raghav B. Venkataramaiyer",
        ],
        14, False, RGBColor(0xB8, 0xC4, 0xD2))
    # no standard footer on title

    # ----- 2 Elevator / overview -----
    s = prs.slides.add_slide(blank)
    rect(s, 0, 0, prs.slide_width, prs.slide_height, PAPER)
    header(s, prs, "Overview", "Three-sentence summary")
    blocks = [
        ("Problem", "Placement rounds and assignment cutoffs are scattered across WhatsApp, Drive folders, and LMS notices."),
        ("Approach", "DeadlineDesk unifies a Placement Track and an Academic Dropbox in one role-based web system."),
        ("Outcome", "Students complete checklists and submissions on time; TAs apply late policy consistently; progress is measurable."),
    ]
    y = Inches(1.55)
    for title, body in blocks:
        rect(s, Inches(0.5), y, Inches(12.3), Inches(1.35), WHITE, LINE)
        rect(s, Inches(0.5), y, Inches(0.12), Inches(1.35), AMBER)
        txt(s, Inches(0.85), y + Inches(0.22), Inches(11.7), Inches(0.35), [title], 16, True, NAVY)
        txt(s, Inches(0.85), y + Inches(0.58), Inches(11.7), Inches(0.6), [body], 15, False, MUTED)
        y += Inches(1.5)
    footer(s, prs, 2, N)

    # ----- 3 Problem -----
    s = prs.slides.add_slide(blank)
    rect(s, 0, 0, prs.slide_width, prs.slide_height, PAPER)
    header(s, prs, "1 · Problem Statement", "Why the project is needed")
    rect(s, Inches(0.5), Inches(1.5), Inches(6.0), Inches(4.7), WHITE, LINE)
    txt(s, Inches(0.75), Inches(1.7), Inches(5.5), Inches(0.4), ["Placement side"], 16, True, NAVY)
    for i, line in enumerate([
        "Company rounds tracked in sheets and chats",
        "Document checklists are informal",
        "Reminders are ad-hoc; windows are missed",
    ]):
        txt(s, Inches(0.75), Inches(2.25) + Inches(i * 0.55), Inches(5.5), Inches(0.5), [f"•  {line}"], 14, False, MUTED)

    rect(s, Inches(6.8), Inches(1.5), Inches(6.0), Inches(4.7), WHITE, LINE)
    txt(s, Inches(7.05), Inches(1.7), Inches(5.5), Inches(0.4), ["Academic side"], 16, True, NAVY)
    for i, line in enumerate([
        "Submissions via Drive or email",
        "Late policy applied manually",
        "Grade and late status lack transparency",
    ]):
        txt(s, Inches(7.05), Inches(2.25) + Inches(i * 0.55), Inches(5.5), Inches(0.5), [f"•  {line}"], 14, False, MUTED)

    txt(s, Inches(0.5), Inches(6.35), Inches(12), Inches(0.35),
        ["Stakeholders: students, teaching assistants / faculty, placement coordinators."], 13, False, MUTED)
    footer(s, prs, 3, N)

    # ----- 4 Objectives -----
    s = prs.slides.add_slide(blank)
    rect(s, 0, 0, prs.slide_width, prs.slide_height, PAPER)
    header(s, prs, "2 · Objectives", "SMART outcomes for the semester")
    txt(s, Inches(0.5), Inches(1.45), Inches(12.3), Inches(0.55),
        ["Primary goal: a working, testable deadline system by Week 7; design depth and quality for EST."],
        15, True, NAVY)
    objs = [
        "O1  Role-based auth for Student, TA/Faculty, Placement Admin (with login tests)",
        "O2  Placement path: company / round → document checklist → T−24h reminder log",
        "O3  Academic path: assignment + late policy → submit → late flag → similarity stub → grade",
        "O4  ≥ 5 automated tests; 100% agreement with late-policy rules on a fixed suite",
        "O5  CI on every push; MkDocs documentation site via GitHub Pages",
    ]
    y = Inches(2.15)
    for o in objs:
        rect(s, Inches(0.5), y, Inches(12.3), Inches(0.68), WHITE, LINE)
        txt(s, Inches(0.75), y + Inches(0.18), Inches(11.9), Inches(0.4), [o], 15, False, INK)
        y += Inches(0.78)
    footer(s, prs, 4, N)

    # ----- 5 Solution -----
    s = prs.slides.add_slide(blank)
    rect(s, 0, 0, prs.slide_width, prs.slide_height, PAPER)
    header(s, prs, "3 · Proposed Solution", "DeadlineDesk — two modules, one core")
    rect(s, Inches(3.6), Inches(1.5), Inches(6.1), Inches(1.1), NAVY)
    txt(s, Inches(3.75), Inches(1.75), Inches(5.8), Inches(0.6),
        ["Shared core: Auth · Roles · Deadlines · Documents · Reminders"], 14, True, WHITE, PP_ALIGN.CENTER)

    rect(s, Inches(0.5), Inches(3.0), Inches(5.9), Inches(3.0), WHITE, LINE)
    rect(s, Inches(0.5), Inches(3.0), Inches(5.9), Inches(0.5), GREEN)
    txt(s, Inches(0.7), Inches(3.1), Inches(5.5), Inches(0.35), ["Placement Track"], 15, True, WHITE)
    for i, line in enumerate(["Companies and rounds", "Document checklists", "Reminder schedule (T−72h / T−24h)"]):
        txt(s, Inches(0.75), Inches(3.7) + Inches(i * 0.45), Inches(5.4), Inches(0.4), [f"•  {line}"], 14, False, MUTED)

    rect(s, Inches(6.9), Inches(3.0), Inches(5.9), Inches(3.0), WHITE, LINE)
    rect(s, Inches(6.9), Inches(3.0), Inches(5.9), Inches(0.5), AMBER)
    txt(s, Inches(7.1), Inches(3.1), Inches(5.5), Inches(0.35), ["Academic Dropbox"], 15, True, WHITE)
    for i, line in enumerate(["Assignments with late policy", "File submit and late flag", "Similarity stub and TA grade"]):
        txt(s, Inches(7.15), Inches(3.7) + Inches(i * 0.45), Inches(5.4), Inches(0.4), [f"•  {line}"], 14, False, MUTED)
    footer(s, prs, 5, N)

    # ----- 6 Methodology -----
    s = prs.slides.add_slide(blank)
    rect(s, 0, 0, prs.slide_width, prs.slide_height, PAPER)
    header(s, prs, "4 · Methodology", "Architecture, stack, and validation")
    table_slide(
        s,
        prs,
        [
            ["Layer", "Responsibility"],
            ["Presentation", "Responsive web UI (student, TA, admin)"],
            ["Application", "Auth, CRUD, status transitions, reminder hooks"],
            ["Domain", "Rounds, checklists, submissions, grades, audit"],
            ["Persistence", "PostgreSQL (SQLite for early prototype)"],
            ["Quality", "Pytest/Jest · GitHub Actions CI · MkDocs Pages"],
        ],
        [Inches(3.2), Inches(9.1)],
    )
    txt(s, Inches(0.5), Inches(4.4), Inches(12.3), Inches(1.5),
        [
            "Stack: Django or Next.js (chosen by delivery speed). Similarity check is an explicit stub.",
            "Primary metric: On-time Completion Rate (OTCR) from database timestamps.",
            "Secondary: late-policy test agreement, reminder effectiveness, TA turnaround.",
            "Distinct from CaseRoom (interview scoring) and College Mentorship (mentor matching).",
        ],
        14, False, MUTED)
    footer(s, prs, 6, N)

    # ----- 7 Timeline -----
    s = prs.slides.add_slide(blank)
    rect(s, 0, 0, prs.slide_width, prs.slide_height, PAPER)
    header(s, prs, "5 · Timeline", "Week-by-week plan to EST")
    table_slide(
        s,
        prs,
        [
            ["Week", "Focus", "Deliverable"],
            ["W4", "Proposal and pitch", "PDF + slides + repository"],
            ["W5–W6", "SRS, diagrams, auth scaffold", "SRS v0.1, CI green"],
            ["W7", "Both end-to-end paths + tests", "Prototype (MST)"],
            ["W8", "Improvement plan", "CE note"],
            ["W11–W12", "UML depth, tests, UI polish", "Design pack"],
            ["W15", "Second prototype", "Prototype 2"],
            ["W17", "Final demo and report", "EST delivery"],
        ],
        [Inches(1.8), Inches(5.2), Inches(5.3)],
    )
    txt(s, Inches(0.5), Inches(5.9), Inches(12), Inches(0.5),
        ["Buffer: W13 and mid-semester weeks absorb delay. Scope cuts prefer dashboards over core status logic."],
        13, False, MUTED)
    footer(s, prs, 7, N)

    # ----- 8 Team / resources -----
    s = prs.slides.add_slide(blank)
    rect(s, 0, 0, prs.slide_width, prs.slide_height, PAPER)
    header(s, prs, "6 · Resources", "Team roles and feasibility")
    table_slide(
        s,
        prs,
        [
            ["Member", "Roll", "Primary responsibility"],
            ["Aryan Gupta", "1024030764", "Architecture, CI/CD, Placement Track"],
            ["Aksh Goyal", "1024030766", "Academic Dropbox, late-policy tests"],
            ["Naveen Bansal", "1024030767", "UI flows, documentation site, demo script"],
        ],
        [Inches(3.5), Inches(2.4), Inches(6.4)],
    )
    txt(s, Inches(0.5), Inches(3.8), Inches(12.3), Inches(2.0),
        [
            "Emails: agupta_be24@thapar.edu · agoyal2_be24@thapar.edu · nbansal3_be24@thapar.edu",
            "Budget: no external funding. GitHub Actions, Pages, and open-source stack only.",
            "No specialised laboratory hardware. Staging on free-tier hosting.",
            "Repository: github.com/aryanorastar/ucs503p-202627-deadlinedesk",
            "Project page: aryanorastar.github.io/ucs503p-202627-deadlinedesk",
        ],
        14, False, MUTED)
    footer(s, prs, 8, N)

    # ----- 9 Risks -----
    s = prs.slides.add_slide(blank)
    rect(s, 0, 0, prs.slide_width, prs.slide_height, PAPER)
    header(s, prs, "7 · Risks", "Anticipated issues and mitigations")
    table_slide(
        s,
        prs,
        [
            ["Risk", "Mitigation"],
            ["Two-module scope creep", "Shared domain core; thin Week-7 paths only"],
            ["Reminder delivery limits", "Persist reminder log first; email later"],
            ["Over-expectation on similarity check", "Document as stub in SRS; TA decides"],
            ["Uneven team velocity", "Rebalance roles after Week-7 review"],
        ],
        [Inches(4.5), Inches(7.8)],
    )
    footer(s, prs, 9, N)

    # ----- 10 Close -----
    s = prs.slides.add_slide(blank)
    rect(s, 0, 0, prs.slide_width, prs.slide_height, NAVY)
    rect(s, Inches(0.5), Inches(2.0), Inches(1.2), Inches(0.08), AMBER)
    txt(s, Inches(0.5), Inches(2.3), Inches(12), Inches(1.0),
        ["Request for approval"], 28, True, WHITE, name="Georgia")
    txt(s, Inches(0.5), Inches(3.4), Inches(12), Inches(2.2),
        [
            "We request registration of DeadlineDesk on the team sheet.",
            "Week-7 commitment: authentication, placement checklist path,",
            "and academic submit / late / grade path with automated tests.",
            "",
            "Thank you. Questions are welcome.",
        ],
        18, False, RGBColor(0xD0, 0xD8, 0xE2))
    txt(s, Inches(0.5), Inches(6.3), Inches(12), Inches(0.6),
        ["Aryan Gupta  ·  Aksh Goyal  ·  Naveen Bansal  ·  UCS503P 2026–27"], 13, False, AMBER)

    out = "w4/DeadlineDesk_W4_Pitch.pptx"
    prs.save(out)
    print("Wrote", out)


if __name__ == "__main__":
    build()
